from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Cm,Pt
import time
import subprocess
import configparser
import os



week     = ['主日','週一', '週二', '週三', '週四', '週五', '週六']

def addSlide(prs,slideText):
	#Add Background
	slide		= prs.slides.add_slide(slideLayout)
	picture     = slide.shapes.add_picture(configBg,Cm(0), Cm(0))
	# This moves it to the background
	slide.shapes._spTree.remove(picture._element)
	slide.shapes._spTree.insert(2, picture._element)


	title		= slide.shapes.title
	title.left  = Cm(1.1)
	title.top  = Cm(1.1)
	title.width  = Cm(32)
	title.height = Cm(18)

	title.text   = ''.join(str(text) for text in slideText)
	for i in range(len(title.text_frame.paragraphs)):		
		title.text_frame.paragraphs[i].alignment = PP_ALIGN.LEFT
		title.text_frame.paragraphs[i].line_spacing = 1.0
	title.text_frame.vertical_anchor = MSO_ANCHOR.TOP
	title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

	paragraphsNum = len(title.text_frame.paragraphs)
	for i in range(paragraphsNum):	
		font = title.text_frame.paragraphs[i].font
		font.name = configFontName
		font.size = configFontSize
		font.bold = configBold

if __name__ == '__main__':

	#Read PPT directory
	prs = Presentation('./template.pptx')
	slideLayout = prs.slide_layouts[5]

	#Read Config File
	config    = configparser.ConfigParser()
	config.sections()
	config.read('./config.ini', encoding = 'UTF-16')
	configBg       = ('./bg.jpg')
	configBold     = bool(config['font']['fontBold'])
	configFontName = config['font']['fontName']
	configFontSize = Pt(int(config['font']['fontSize']))

	#Read Input File
	textFile		= open('./input.txt', 'r', encoding = 'UTF-8')
	inputText	   = textFile.readlines()
	if(inputText[-1] != '\n'):
		inputText.append('\n')
	slideText	   = []

	#Read Input line
	titleLine = inputText[0][:-1]+".pptx"
	print(titleLine)
	for line in inputText:
		line = line.lstrip(" ")
		if (line == '\n'):
			addSlide(prs,slideText)
			slideText = []
		else:
			slideText.append(line)

	#Set FileName
	dateType = time.strftime('%Y%m%d', time.localtime()) + week[int(time.strftime('%w', time.localtime()))]
	outputName = dateType + '清晨箴言.pptx'

	#Output File
	print("箴言投影片完成")
	prs.save('./output/'+ str(titleLine))
	subprocess.Popen(os.path.abspath("output/" + str(titleLine)), shell=True)